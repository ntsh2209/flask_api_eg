from pptx import Presentation
from pptx.util import Inches
import pandas as pd

def find_slide_by_title(presentation, slide_title):
    """Find a slide by its title."""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text == slide_title:
                return slide
    raise ValueError(f"Slide with title '{slide_title}' not found")

def add_dataframe_to_slide(slide, df):
    """Add DataFrame content to a slide as a table."""
    # Define the position and size of the table
    rows, cols = df.shape
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8 * rows)
    
    # Add table to slide
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    
    # Add column names (header row)
    for col_idx, col_name in enumerate(df.columns):
        table.cell(0, col_idx).text = col_name
    
    # Add DataFrame content to table
    for row_idx, row in df.iterrows():
        for col_idx, value in enumerate(row):
            table.cell(row_idx + 1, col_idx).text = str(value)

def modify_ppt_with_content(ppt_file, slide_title, df):
    """Modify an existing PowerPoint presentation and add content from DataFrame."""
    # Load the PowerPoint presentation
    presentation = Presentation(ppt_file)
    
    # Find the slide by its title
    slide = find_slide_by_title(presentation, slide_title)
    
    # Add the DataFrame as a table to the slide
    add_dataframe_to_slide(slide, df)
    
    # Save the modified presentation
    new_ppt_file = f"modified_{ppt_file}"
    presentation.save(new_ppt_file)
    print(f"Presentation saved as: {new_ppt_file}")

# Example usage
if __name__ == "__main__":
    # Load your DataFrame (example)
    data = {'Name': ['Alice', 'Bob', 'Charlie'],
            'Age': [25, 30, 35],
            'City': ['New York', 'Los Angeles', 'Chicago']}
    df = pd.DataFrame(data)
    
    # File path to the existing PowerPoint
    ppt_file = 'your_presentation.pptx'
    
    # Slide title where the content should be added
    slide_title = 'Your Slide Title'  # Title of the slide to modify
    
    # Modify the PowerPoint and add the DataFrame content
    modify_ppt_with_content(ppt_file, slide_title, df)
